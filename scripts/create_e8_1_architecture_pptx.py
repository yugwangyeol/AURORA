#!/usr/bin/env python3
"""Create a one-slide, editable E8.1 architecture figure."""

from __future__ import annotations

import tempfile
from pathlib import Path

import pymupdf
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path("/home/jovyan")
REFERENCE_PDF = ROOT / "main_fig.pdf"
OUTPUT_PPTX = ROOT / "E8_1_architecture_figure.pptx"


def rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


C = {
    "bg": rgb("F8FAFC"),
    "ink": rgb("172B3A"),
    "muted": rgb("64748B"),
    "line": rgb("94A3B8"),
    "panel": rgb("EFF6FF"),
    "panel_line": rgb("B6CAE8"),
    "teal": rgb("0F766E"),
    "teal_soft": rgb("CCFBF1"),
    "orange": rgb("EA580C"),
    "orange_soft": rgb("FFEDD5"),
    "blue": rgb("2563EB"),
    "blue_soft": rgb("DBEAFE"),
    "purple": rgb("7C3AED"),
    "purple_soft": rgb("EDE9FE"),
    "yellow": rgb("B77900"),
    "yellow_soft": rgb("FEF3C7"),
    "red": rgb("C62828"),
    "red_soft": rgb("FEE2E2"),
    "white": rgb("FFFFFF"),
    "gray_soft": rgb("F1F5F9"),
}


def set_text(
    shape,
    text: str,
    *,
    size: float = 13,
    color: RGBColor | None = None,
    bold: bool = False,
    align=PP_ALIGN.CENTER,
    font: str = "Arial",
    margin: float = 0.06,
) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or C["ink"]


def textbox(slide, x, y, w, h, text, **kwargs):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(shape, text, **kwargs)
    return shape


def box(
    slide,
    x,
    y,
    w,
    h,
    text="",
    *,
    fill=None,
    line=None,
    radius=True,
    size=13,
    color=None,
    bold=False,
    width=1.2,
    align=PP_ALIGN.CENTER,
    font="Arial",
):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill or C["white"]
    shape.line.color.rgb = line or C["line"]
    shape.line.width = Pt(width)
    set_text(
        shape,
        text,
        size=size,
        color=color,
        bold=bold,
        align=align,
        font=font,
    )
    return shape


def line(
    slide,
    x1,
    y1,
    x2,
    y2,
    *,
    color=None,
    width=1.5,
    dashed=False,
    arrow=True,
):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    connector.line.color.rgb = color or C["line"]
    connector.line.width = Pt(width)
    if dashed:
        connector.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if arrow:
        ln = connector.line._get_or_add_ln()
        tail = OxmlElement("a:tailEnd")
        tail.set("type", "triangle")
        tail.set("w", "sm")
        tail.set("len", "sm")
        ln.append(tail)
    return connector


def extract_reference_photo(destination: Path) -> None:
    doc = pymupdf.open(REFERENCE_PDF)
    candidates = []
    for image in doc[0].get_images(full=True):
        xref, width, height = image[0], image[2], image[3]
        candidates.append((width * height, xref))
    _, xref = max(candidates)
    extracted = doc.extract_image(xref)
    destination.write_bytes(extracted["image"])


def add_patch_lane(slide, x: float, y: float, count: int = 16) -> None:
    colors = ["DBEAFE", "BFDBFE", "D1FAE5", "E0E7FF"]
    for i in range(count):
        px = x + i * 0.365
        patch = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(px), Inches(y), Inches(0.29), Inches(0.24)
        )
        patch.fill.solid()
        patch.fill.fore_color.rgb = rgb(colors[i % len(colors)])
        patch.line.color.rgb = C["panel_line"]
        patch.line.width = Pt(0.6)


def add_mask_icon(slide, x: float, y: float, scale: float = 1.0) -> None:
    mask = [
        0, 0, 0, 0,
        0, 1, 1, 0,
        0, 1, 2, 2,
        0, 0, 2, 2,
    ]
    colors = {0: C["gray_soft"], 1: C["teal_soft"], 2: C["orange_soft"]}
    outlines = {0: C["line"], 1: C["teal"], 2: C["orange"]}
    side = 0.125 * scale
    for r in range(4):
        for c in range(4):
            value = mask[r * 4 + c]
            cell = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x + c * side),
                Inches(y + r * side),
                Inches(side),
                Inches(side),
            )
            cell.fill.solid()
            cell.fill.fore_color.rgb = colors[value]
            cell.line.color.rgb = outlines[value]
            cell.line.width = Pt(0.35)


def add_writer(slide, x: float, layer_idx: int, memory_label: str) -> None:
    textbox(
        slide,
        x + 0.30,
        1.27,
        1.30,
        0.25,
        f"LAYER {layer_idx}",
        size=9,
        color=C["blue"],
        bold=True,
    )
    box(
        slide,
        x + 0.10,
        1.58,
        1.70,
        0.54,
        f"semantic key  sₖ^{layer_idx}",
        fill=C["teal_soft"],
        line=C["teal"],
        color=C["teal"],
        size=11,
        bold=True,
    )
    line(slide, x + 0.95, 2.12, x + 0.95, 2.55, color=C["teal"], width=1.5)
    box(
        slide,
        x,
        2.55,
        1.90,
        1.55,
        "COMPETITIVE WRITER\nQ: semantic + memory\nK,V: image patches\nα = softmax_slots(QKᵀ)\nmₖ ← Σₚ ᾱₖₚ V(xₚ)",
        fill=C["white"],
        line=C["blue"],
        color=C["ink"],
        size=8.8,
        bold=False,
        width=1.6,
        font="Cambria Math",
    )
    line(slide, x + 0.95, 4.10, x + 0.95, 4.43, color=C["orange"], width=1.7)
    box(
        slide,
        x + 0.10,
        4.43,
        1.70,
        0.55,
        memory_label,
        fill=C["orange_soft"],
        line=C["orange"],
        color=C["orange"],
        size=11,
        bold=True,
        font="Cambria Math",
    )
    add_mask_icon(slide, x + 1.57, 2.67, scale=0.48)
    textbox(slide, x + 1.30, 2.92, 0.54, 0.18, "αᶫ", size=7.5, color=C["red"], bold=True)


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = C["bg"]

    # Header
    textbox(
        slide,
        0.35,
        0.14,
        10.8,
        0.48,
        "E8.1  ·  Semantic-Keyed Object Visual Memory inside an MLLM",
        size=22,
        color=C["ink"],
        bold=True,
        align=PP_ALIGN.LEFT,
        margin=0,
    )
    textbox(
        slide,
        0.37,
        0.62,
        9.8,
        0.25,
        "Object identity chooses where to read; image-only memory supplies what to reconstruct.",
        size=11.5,
        color=C["muted"],
        align=PP_ALIGN.LEFT,
        margin=0,
    )
    box(
        slide,
        12.28,
        0.18,
        3.34,
        0.52,
        "VARIABLE K OBJECTS  +  4 BG REGISTERS",
        fill=C["gray_soft"],
        line=C["panel_line"],
        color=C["muted"],
        size=10.5,
        bold=True,
    )

    with tempfile.TemporaryDirectory() as tmp:
        photo = Path(tmp) / "input.jpg"
        extract_reference_photo(photo)
        slide.shapes.add_picture(str(photo), Inches(0.35), Inches(2.16), Inches(1.38), Inches(1.38))
    border = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.35), Inches(2.16), Inches(1.38), Inches(1.38)
    )
    border.fill.background()
    border.line.color.rgb = C["line"]
    border.line.width = Pt(1.0)
    textbox(slide, 0.35, 3.56, 1.38, 0.25, "INPUT IMAGE", size=10, bold=True)

    # Input encoders and prompt
    line(slide, 1.73, 2.85, 1.96, 2.85, color=C["line"], width=1.8)
    box(
        slide,
        1.96,
        2.30,
        1.32,
        1.10,
        "FROZEN\nDUAL-SCALE\nVISION ENCODER",
        fill=C["gray_soft"],
        line=C["line"],
        size=10.5,
        bold=True,
    )
    box(
        slide,
        1.72,
        1.05,
        1.78,
        0.88,
        "OBJECT PROMPT\ncaptionₖ + <OVT>ₖ",
        fill=C["teal_soft"],
        line=C["teal"],
        color=C["teal"],
        size=11,
        bold=True,
        font="Cambria Math",
    )
    line(slide, 3.02, 1.93, 3.68, 1.76, color=C["teal"], width=1.7)
    textbox(slide, 1.98, 3.47, 1.25, 0.30, "image patches  xₚ", size=9.5, color=C["blue"], font="Cambria Math")
    line(slide, 3.28, 2.85, 3.63, 2.85, color=C["blue"], width=1.8)
    box(
        slide,
        0.55,
        4.45,
        1.75,
        0.62,
        "4 LEARNABLE\nBG REGISTERS  rⱼ",
        fill=C["purple_soft"],
        line=C["purple"],
        color=C["purple"],
        size=10.5,
        bold=True,
        font="Cambria Math",
    )
    line(slide, 2.30, 4.76, 3.68, 4.76, color=C["purple"], width=1.6)

    # MLLM container
    panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.58),
        Inches(0.92),
        Inches(7.16),
        Inches(5.32),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = C["panel"]
    panel.line.color.rgb = C["panel_line"]
    panel.line.width = Pt(1.4)
    textbox(
        slide,
        3.82,
        0.99,
        5.0,
        0.34,
        "QWEN MLLM  ·  ITERATIVE VISUAL WRITE",
        size=12.5,
        color=C["blue"],
        bold=True,
        align=PP_ALIGN.LEFT,
        margin=0,
    )
    textbox(slide, 8.00, 1.02, 2.45, 0.28, "shared Writer parameters", size=9.5, color=C["muted"], align=PP_ALIGN.RIGHT)
    textbox(slide, 3.78, 1.38, 2.2, 0.22, "SEMANTIC STREAM", size=8.5, color=C["teal"], bold=True, align=PP_ALIGN.LEFT)
    textbox(slide, 3.78, 2.08, 2.25, 0.22, "IMAGE PATCH STREAM", size=8.5, color=C["blue"], bold=True, align=PP_ALIGN.LEFT)
    textbox(slide, 3.78, 4.22, 2.4, 0.22, "VISUAL MEMORY STREAM", size=8.5, color=C["orange"], bold=True, align=PP_ALIGN.LEFT)
    add_patch_lane(slide, 4.03, 2.31, 17)

    writer_xs = [3.92, 6.20, 8.48]
    for x, layer_idx, memory_label in zip(
        writer_xs,
        [21, 24, 27],
        ["visual memory  mₖ²¹", "refined memory  mₖ²⁴", "final memory  mₖ²⁷"],
    ):
        add_writer(slide, x, layer_idx, memory_label)

    line(slide, 5.72, 1.85, 6.30, 1.85, color=C["teal"], width=1.4)
    line(slide, 8.00, 1.85, 8.58, 1.85, color=C["teal"], width=1.4)
    line(slide, 5.72, 4.71, 6.30, 4.71, color=C["orange"], width=1.7)
    line(slide, 8.00, 4.71, 8.58, 4.71, color=C["orange"], width=1.7)
    for x in [4.87, 7.15, 9.43]:
        line(slide, x, 2.54, x, 2.64, color=C["blue"], width=1.0)

    box(
        slide,
        4.08,
        5.36,
        6.16,
        0.48,
        "CLEAN BOTTLENECK: visual memory is NOT injected back into the Qwen residual stream",
        fill=C["white"],
        line=C["orange"],
        color=C["orange"],
        size=10.2,
        bold=True,
    )

    # Typed Reader
    line(slide, 10.06, 1.85, 11.05, 2.42, color=C["teal"], width=2.0)
    textbox(slide, 10.48, 1.82, 0.65, 0.25, "KEYS", size=8.5, color=C["teal"], bold=True)
    line(slide, 10.06, 4.75, 11.05, 4.04, color=C["orange"], width=2.2)
    textbox(slide, 10.42, 4.43, 0.75, 0.25, "VALUES", size=8.5, color=C["orange"], bold=True)
    box(
        slide,
        11.05,
        2.08,
        2.32,
        2.58,
        "Aᵢₖ = softmaxₖ(Q(qᵢ)ᵀK(sₖ))\n\ncᵢ = Σₖ Aᵢₖ V(mₖ)",
        fill=C["blue_soft"],
        line=C["blue"],
        color=C["ink"],
        size=11.5,
        bold=False,
        width=1.8,
        font="Cambria Math",
    )
    textbox(slide, 11.23, 2.18, 1.96, 0.28, "TYPED RAE READER", size=12.5, color=C["blue"], bold=True)
    box(
        slide,
        11.38,
        5.04,
        1.65,
        0.57,
        "learnable RAE queries  qᵢ",
        fill=C["purple_soft"],
        line=C["purple"],
        color=C["purple"],
        size=10,
        bold=True,
        font="Cambria Math",
    )
    line(slide, 12.20, 5.04, 12.20, 4.67, color=C["purple"], width=1.6)
    textbox(slide, 11.15, 4.60, 2.12, 0.30, "values come ONLY from mₖ", size=9, color=C["orange"], bold=True)

    # DiT
    line(slide, 13.37, 3.25, 13.74, 3.25, color=C["blue"], width=2.2)
    textbox(slide, 13.38, 2.94, 0.42, 0.22, "cᵢ", size=10, color=C["blue"], bold=True, font="Cambria Math")
    box(
        slide,
        13.74,
        2.34,
        1.88,
        1.82,
        "DIFFUSION\nTRANSFORMER\n(DiT)\n\nCross-attn + AdaLN",
        fill=C["yellow_soft"],
        line=C["yellow"],
        color=C["yellow"],
        size=12,
        bold=True,
        width=1.8,
    )
    box(
        slide,
        14.10,
        1.35,
        1.18,
        0.54,
        "noise  zₜ",
        fill=C["gray_soft"],
        line=C["line"],
        color=C["muted"],
        size=10.5,
        bold=True,
        font="Cambria Math",
    )
    line(slide, 14.69, 1.89, 14.69, 2.34, color=C["line"], width=1.5)

    # Losses
    box(
        slide,
        11.18,
        0.91,
        2.02,
        0.66,
        "Caption logits  ↔  caption labels\nL_LM  (autoregressive CE)",
        fill=C["red_soft"],
        line=C["red"],
        color=C["red"],
        size=9.5,
        bold=True,
        font="Cambria Math",
    )
    line(slide, 10.74, 1.28, 11.18, 1.28, color=C["red"], width=1.3, dashed=True)

    gt_box = box(
        slide,
        0.35,
        6.38,
        2.15,
        0.82,
        "",
        fill=C["red_soft"],
        line=C["red"],
        color=C["red"],
    )
    add_mask_icon(slide, 0.50, 6.53, scale=0.78)
    textbox(
        slide,
        1.00,
        6.46,
        1.38,
        0.62,
        "GT INSTANCE MASKS\nM_GT · object + bg",
        size=9.5,
        color=C["red"],
        bold=True,
        align=PP_ALIGN.LEFT,
        font="Cambria Math",
    )
    line(slide, 2.50, 6.79, 4.05, 6.79, color=C["red"], width=1.4, dashed=True)
    box(
        slide,
        4.05,
        6.34,
        2.62,
        0.90,
        "WRITER OWNERSHIP LOSS\nL_owner = ⅓ Σᶫ CE(αᶫ, M_GT)\nℓ ∈ {21, 24, 27}",
        fill=C["white"],
        line=C["red"],
        color=C["red"],
        size=10.2,
        bold=True,
        font="Cambria Math",
    )
    for x in [5.62, 7.90, 10.18]:
        line(slide, x, 3.20, x, 6.12, color=C["red"], width=0.9, dashed=True, arrow=False)
    line(slide, 5.62, 6.12, 10.18, 6.12, color=C["red"], width=0.9, dashed=True, arrow=False)
    line(slide, 5.36, 6.12, 5.36, 6.34, color=C["red"], width=1.1, dashed=True)

    line(slide, 2.50, 7.05, 10.64, 7.05, color=C["red"], width=1.0, dashed=True)
    line(slide, 13.37, 4.12, 13.42, 6.14, color=C["red"], width=1.2, dashed=True, arrow=False)
    line(slide, 13.42, 6.14, 11.88, 6.34, color=C["red"], width=1.2, dashed=True)
    box(
        slide,
        10.64,
        6.34,
        2.48,
        0.90,
        "READER ROUTING LOSS\nL_reader = CE(A, ↓₁₆ M_GT)\nobject + aggregate background",
        fill=C["white"],
        line=C["red"],
        color=C["red"],
        size=10,
        bold=True,
        font="Cambria Math",
    )

    box(
        slide,
        0.55,
        5.31,
        2.35,
        0.62,
        "target image → frozen RAE latent  v",
        fill=C["gray_soft"],
        line=C["line"],
        color=C["muted"],
        size=10,
        bold=True,
        font="Cambria Math",
    )
    line(slide, 1.04, 3.81, 1.04, 5.31, color=C["line"], width=1.2)
    line(slide, 2.90, 5.62, 2.90, 5.98, color=C["line"], width=1.1, dashed=True, arrow=False)
    line(slide, 2.90, 5.98, 14.67, 5.98, color=C["line"], width=1.1, dashed=True)
    line(slide, 14.67, 4.16, 14.67, 6.31, color=C["yellow"], width=1.5)
    box(
        slide,
        13.52,
        6.31,
        2.12,
        0.93,
        "FLOW-MATCHING LOSS\nL_recon = FM(DiT(c, zₜ), v)",
        fill=C["white"],
        line=C["red"],
        color=C["red"],
        size=10.5,
        bold=True,
        font="Cambria Math",
    )

    # Legend and total objective
    box(slide, 0.35, 7.51, 0.32, 0.28, "", fill=C["teal_soft"], line=C["teal"], radius=False)
    textbox(slide, 0.72, 7.48, 1.45, 0.32, "semantic key / identity", size=8.5, color=C["muted"], align=PP_ALIGN.LEFT)
    box(slide, 2.25, 7.51, 0.32, 0.28, "", fill=C["orange_soft"], line=C["orange"], radius=False)
    textbox(slide, 2.62, 7.48, 1.55, 0.32, "image-only appearance", size=8.5, color=C["muted"], align=PP_ALIGN.LEFT)
    box(slide, 4.25, 7.51, 0.32, 0.28, "", fill=C["purple_soft"], line=C["purple"], radius=False)
    textbox(slide, 4.62, 7.48, 1.45, 0.32, "query / background", size=8.5, color=C["muted"], align=PP_ALIGN.LEFT)
    line(slide, 6.28, 7.65, 6.78, 7.65, color=C["red"], width=1.1, dashed=True, arrow=False)
    textbox(slide, 6.82, 7.48, 1.30, 0.32, "supervision", size=8.5, color=C["muted"], align=PP_ALIGN.LEFT)

    box(
        slide,
        0.35,
        8.08,
        15.28,
        0.56,
        "E8.1 OBJECTIVE   L_total = 1.0 L_LM + 1.5 L_recon + 1.0 L_owner + 0.5 L_reader,obj + 0.25 L_reader,bg",
        fill=C["ink"],
        line=C["ink"],
        color=C["white"],
        size=13,
        bold=True,
        font="Cambria Math",
    )

    prs.save(OUTPUT_PPTX)
    print(OUTPUT_PPTX)


if __name__ == "__main__":
    build()
