#!/usr/bin/env python3
"""Create a minimal, paper-style E8.1 architecture figure."""

from __future__ import annotations

import shutil
import tempfile
from pathlib import Path

import pymupdf
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path("/home/jovyan")
REFERENCE_PDF = ROOT / "main_fig.pdf"
OUTPUT_PPTX = ROOT / "E8_1_architecture_figure.pptx"


def rgb(code: str) -> RGBColor:
    code = code.lstrip("#")
    return RGBColor(*(int(code[i : i + 2], 16) for i in (0, 2, 4)))


INK = rgb("162B3A")
MUTED = rgb("64748B")
LINE = rgb("9AA5B1")
MLLM_BG = rgb("EEF4FF")
MLLM_LINE = rgb("BED0EE")
TEAL = rgb("16877E")
TEAL_BG = rgb("D8F4EF")
ORANGE = rgb("E76F25")
ORANGE_BG = rgb("FDE8D8")
BLUE = rgb("3478C7")
BLUE_BG = rgb("DCEBFA")
PURPLE = rgb("8064A2")
PURPLE_BG = rgb("E9E1F2")
YELLOW = rgb("A97400")
YELLOW_BG = rgb("FFF1C7")
RED = rgb("C62828")
RED_BG = rgb("FDE6E6")
GRAY_BG = rgb("EFF1F3")
WHITE = rgb("FFFFFF")


def set_text(shape, text, *, size=12, color=INK, bold=False, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    r = p.add_run()
    r.text = text
    r.font.name = "Arial"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def text(slide, x, y, w, h, value, **kwargs):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(shape, value, **kwargs)
    return shape


def box(slide, x, y, w, h, value="", *, fill=WHITE, stroke=LINE, size=12, color=INK, bold=False, kind=MSO_SHAPE.ROUNDED_RECTANGLE, width=1.3):
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = stroke
    shape.line.width = Pt(width)
    set_text(shape, value, size=size, color=color, bold=bold)
    return shape


def arrow(slide, x1, y1, x2, y2, *, color=LINE, width=1.6, dashed=False, head=True):
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    if dashed:
        shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if head:
        tail = OxmlElement("a:tailEnd")
        tail.set("type", "triangle")
        tail.set("w", "sm")
        tail.set("len", "sm")
        shape.line._get_or_add_ln().append(tail)
    return shape


def loss(slide, x, y, w, label):
    return box(
        slide,
        x,
        y,
        w,
        0.42,
        label,
        fill=RED_BG,
        stroke=RED,
        color=RED,
        size=11.5,
        bold=True,
        width=1.2,
    )


def token(slide, x, y, w, label, fill, stroke, *, size=11):
    return box(
        slide,
        x,
        y,
        w,
        0.46,
        label,
        fill=fill,
        stroke=stroke,
        color=stroke,
        size=size,
        bold=True,
        width=1.3,
    )


def slot_pair(slide, x, y, semantic_label, register_label):
    token(slide, x, y, 0.78, semantic_label, TEAL_BG, TEAL)
    token(slide, x + 0.84, y, 0.44, register_label, PURPLE_BG, PURPLE, size=10)


def memory_pair(slide, x, y, object_label, bg_label):
    token(slide, x, y, 0.78, object_label, ORANGE_BG, ORANGE)
    token(slide, x + 0.84, y, 0.44, bg_label, PURPLE_BG, PURPLE, size=9)


def mask_icon(slide, x, y, scale=1.0):
    pattern = [0, 0, 0, 0, 0, 1, 1, 0, 0, 1, 2, 2, 0, 0, 2, 2]
    fills = [GRAY_BG, TEAL_BG, ORANGE_BG]
    strokes = [LINE, TEAL, ORANGE]
    side = 0.10 * scale
    for row in range(4):
        for col in range(4):
            v = pattern[row * 4 + col]
            cell = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x + col * side),
                Inches(y + row * side),
                Inches(side),
                Inches(side),
            )
            cell.fill.solid()
            cell.fill.fore_color.rgb = fills[v]
            cell.line.color.rgb = strokes[v]
            cell.line.width = Pt(0.35)


def patches(slide, x, y, count=18):
    palette = [BLUE_BG, rgb("C7DDF5"), TEAL_BG, rgb("E4E9F6")]
    for i in range(count):
        cell = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x + i * 0.325),
            Inches(y),
            Inches(0.25),
            Inches(0.22),
        )
        cell.fill.solid()
        cell.fill.fore_color.rgb = palette[i % len(palette)]
        cell.line.color.rgb = MLLM_LINE
        cell.line.width = Pt(0.4)


def reference_photo(destination: Path):
    doc = pymupdf.open(REFERENCE_PDF)
    candidates = [(img[2] * img[3], img[0]) for img in doc[0].get_images(full=True)]
    _, xref = max(candidates)
    image = doc.extract_image(xref)
    destination.write_bytes(image["image"])


def writer_stage(slide, x, layer, semantic, memory):
    slot_pair(slide, x, 1.18, semantic, "rⱼ")
    text(slide, x + 0.29, 1.76, 0.72, 0.20, f"L{layer}", size=9, color=BLUE, bold=True)
    writer = box(
        slide,
        x + 0.04,
        2.22,
        1.20,
        0.82,
        "Writer\nQ: sₖ + mₖ\nK,V: xₚ",
        fill=WHITE,
        stroke=BLUE,
        color=INK,
        size=8.8,
        bold=True,
        kind=MSO_SHAPE.HEXAGON,
        width=1.5,
    )
    mask_icon(slide, x + 0.97, 2.10, 0.58)
    memory_pair(slide, x, 3.68, memory, "mᵦ")
    arrow(slide, x + 0.64, 1.64, x + 0.64, 2.22, color=TEAL, width=1.5)
    arrow(slide, x + 0.64, 3.04, x + 0.64, 3.68, color=ORANGE, width=1.7)
    return writer


def build():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(6.35)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # Input and image encoding.
    with tempfile.TemporaryDirectory() as tmp:
        photo = Path(tmp) / "input.jpg"
        reference_photo(photo)
        slide.shapes.add_picture(str(photo), Inches(0.24), Inches(2.02), Inches(1.30), Inches(1.30))
    frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.24), Inches(2.02), Inches(1.30), Inches(1.30))
    frame.fill.background()
    frame.line.color.rgb = LINE
    frame.line.width = Pt(0.8)
    text(slide, 0.24, 3.38, 1.30, 0.25, "Input image", size=10.5, bold=True)
    arrow(slide, 1.54, 2.67, 1.88, 2.67, color=LINE, width=1.6)
    box(slide, 1.88, 2.18, 1.35, 0.98, "Dual-scale\nRAE encoder", fill=GRAY_BG, stroke=LINE, size=11, bold=True)
    box(slide, 1.90, 0.68, 1.34, 0.70, "object captionₖ\n+  <OVT>ₖ", fill=TEAL_BG, stroke=TEAL, color=TEAL, size=10.5, bold=True)
    box(slide, 2.12, 4.18, 0.94, 0.52, "registers  rⱼ", fill=PURPLE_BG, stroke=PURPLE, color=PURPLE, size=10, bold=True)

    # Main MLLM block.
    panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.58),
        Inches(0.58),
        Inches(7.10),
        Inches(4.78),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = MLLM_BG
    panel.line.color.rgb = MLLM_LINE
    panel.line.width = Pt(1.4)
    text(slide, 9.18, 0.70, 1.20, 0.28, "Qwen MLLM", size=11, color=BLUE, bold=True)
    patches(slide, 3.92, 1.94, 20)
    text(slide, 3.78, 1.72, 0.65, 0.20, "xₚ", size=10, color=BLUE, bold=True)
    arrow(slide, 3.23, 2.67, 3.65, 2.07, color=BLUE, width=1.8)
    arrow(slide, 3.24, 1.03, 4.00, 1.30, color=TEAL, width=1.7)
    arrow(slide, 3.06, 4.44, 3.86, 4.00, color=PURPLE, width=1.6)

    xs = [4.05, 6.20, 8.35]
    writer_stage(slide, xs[0], 21, "sₖ²¹", "mₖ²¹")
    writer_stage(slide, xs[1], 24, "sₖ²⁴", "mₖ²⁴")
    writer_stage(slide, xs[2], 27, "sₖ²⁷", "mₖ²⁷")

    # Semantic refinement and visual-memory refinement.
    arrow(slide, 5.33, 1.41, 6.20, 1.41, color=TEAL, width=1.6)
    arrow(slide, 7.48, 1.41, 8.35, 1.41, color=TEAL, width=1.6)
    arrow(slide, 5.33, 3.91, 6.20, 3.91, color=ORANGE, width=1.8)
    arrow(slide, 7.48, 3.91, 8.35, 3.91, color=ORANGE, width=1.8)
    arrow(slide, 5.05, 3.68, 6.43, 2.96, color=ORANGE, width=1.2)
    arrow(slide, 7.20, 3.68, 8.58, 2.96, color=ORANGE, width=1.2)
    for center in [4.69, 6.84, 8.99]:
        arrow(slide, center, 2.16, center, 2.23, color=BLUE, width=1.0)

    # Owner and LM losses: names only.
    loss(slide, 6.22, 0.05, 1.12, "L_owner")
    for x in [5.17, 7.32, 9.47]:
        arrow(slide, x, 2.35, x, 0.52, color=RED, width=0.9, dashed=True, head=False)
    arrow(slide, 5.17, 0.52, 9.47, 0.52, color=RED, width=0.9, dashed=True, head=False)
    arrow(slide, 6.78, 0.52, 6.78, 0.46, color=RED, width=1.0, dashed=True)
    loss(slide, 9.30, 0.05, 0.92, "L_LM")
    arrow(slide, 9.74, 0.58, 9.74, 0.47, color=RED, width=1.0, dashed=True)

    # Reader: Q from RAE queries, K from semantics, V from visual memory.
    reader = box(slide, 11.02, 1.32, 2.15, 2.45, "", fill=BLUE_BG, stroke=BLUE, width=1.8)
    text(slide, 11.45, 1.48, 1.28, 0.32, "Reader", size=16, color=BLUE, bold=True)
    token(slide, 11.30, 1.98, 0.46, "Q", PURPLE_BG, PURPLE, size=11)
    text(slide, 11.82, 1.98, 0.94, 0.44, "qᵢ", size=13, color=PURPLE, bold=True)
    token(slide, 11.30, 2.50, 0.46, "K", TEAL_BG, TEAL, size=11)
    text(slide, 11.82, 2.50, 0.94, 0.44, "sₖ , rⱼ", size=13, color=TEAL, bold=True)
    token(slide, 11.30, 3.02, 0.46, "V", ORANGE_BG, ORANGE, size=11)
    text(slide, 11.82, 3.02, 1.02, 0.44, "mₖ , mᵦ", size=13, color=ORANGE, bold=True)
    arrow(slide, 9.63, 1.41, 11.30, 2.72, color=TEAL, width=2.0)
    arrow(slide, 9.63, 3.91, 11.30, 3.24, color=ORANGE, width=2.2)
    token(slide, 11.48, 4.25, 1.28, "RAE queries  qᵢ", PURPLE_BG, PURPLE, size=10.5)
    arrow(slide, 12.12, 4.25, 12.12, 2.42, color=PURPLE, width=1.5)
    mask_icon(slide, 12.54, 3.31, 0.65)
    loss(slide, 11.55, 5.02, 1.18, "L_reader")
    arrow(slide, 12.80, 3.56, 12.20, 5.02, color=RED, width=1.0, dashed=True)

    # DiT reconstruction path.
    token(slide, 13.38, 2.29, 0.50, "cᵢ", BLUE_BG, BLUE, size=11)
    arrow(slide, 13.17, 2.54, 13.38, 2.54, color=BLUE, width=2.0)
    arrow(slide, 13.88, 2.54, 14.05, 2.54, color=BLUE, width=2.0)
    box(slide, 14.05, 1.62, 1.70, 1.84, "Diffusion\nTransformer\n(DiT)", fill=YELLOW_BG, stroke=YELLOW, color=YELLOW, size=12, bold=True, width=1.8)
    token(slide, 14.55, 0.72, 0.82, "noise  zₜ", GRAY_BG, LINE, size=10)
    arrow(slide, 14.96, 1.18, 14.96, 1.62, color=LINE, width=1.4)
    loss(slide, 14.42, 4.28, 1.06, "L_recon")
    arrow(slide, 14.96, 3.46, 14.96, 4.28, color=RED, width=1.2)

    # Frozen reconstruction target path.
    token(slide, 1.98, 3.42, 1.14, "target latent  v", GRAY_BG, LINE, size=10)
    arrow(slide, 2.55, 3.16, 2.55, 3.42, color=LINE, width=1.1)
    arrow(slide, 3.12, 3.65, 3.35, 5.50, color=LINE, width=1.0, dashed=True, head=False)
    arrow(slide, 3.35, 5.50, 14.95, 5.50, color=LINE, width=1.0, dashed=True, head=False)
    arrow(slide, 14.95, 5.50, 14.95, 4.70, color=LINE, width=1.0, dashed=True)

    # Minimal legend.
    token(slide, 0.35, 5.62, 0.52, "sₖ", TEAL_BG, TEAL, size=10)
    text(slide, 0.91, 5.62, 0.78, 0.44, "semantic", size=9.5, color=MUTED)
    token(slide, 1.73, 5.62, 0.52, "mₖ", ORANGE_BG, ORANGE, size=10)
    text(slide, 2.29, 5.62, 0.78, 0.44, "visual", size=9.5, color=MUTED)
    token(slide, 3.11, 5.62, 0.52, "qᵢ", PURPLE_BG, PURPLE, size=10)
    text(slide, 3.67, 5.62, 0.72, 0.44, "query", size=9.5, color=MUTED)

    prs.save(OUTPUT_PPTX)
    print(OUTPUT_PPTX)


if __name__ == "__main__":
    build()
